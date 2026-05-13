import google.generativeai as genai
from pptx import Presentation
from dotenv import load_dotenv
import json
import os
import re
import streamlit as st

load_dotenv()

try:
    api_key = st.secrets["GOOGLE_API_KEY"]
except Exception:
    api_key = os.getenv("GOOGLE_API_KEY")

model = None
if api_key:
    genai.configure(api_key=api_key)
    model = genai.GenerativeModel("gemini-1.5-flash")

TEMPLATES = {
    "pi planning": "templates/pi_planning.pptx",
    "sprint review": "templates/sprint_review.pptx",
    "retrospective": "templates/retrospective.pptx",
    "daily status": "templates/daily_status.pptx",
    "system demo": "templates/System_Demo.pptx",
}

def get_template(user_input):
    user_input = user_input.lower()
    for key, path in TEMPLATES.items():
        if key in user_input:
            return path, key
    return None, None

def clean_json_response(raw_text):
    raw_text = raw_text.strip()
    raw_text = re.sub(r"```json", "", raw_text)
    raw_text = re.sub(r"```", "", raw_text)
    return raw_text.strip()

def extract_content_with_llm(user_input, template_name):
    if not model:
        raise Exception("GOOGLE_API_KEY not configured. Add it in Streamlit Secrets.")

    prompt = f"""
You are an enterprise Agile Release Train presentation assistant.

TEMPLATE SELECTED: {template_name}
USER REQUEST: {user_input}

Return ONLY valid JSON with these fields:
{{
    "title": "Presentation title",
    "objective": "Business objective",
    "agenda": "Agenda items separated by commas",
    "summary": "Executive summary",
    "highlights": "Key achievements or accomplishments",
    "risks": "Major risks or blockers",
    "dependencies": "Dependencies or coordination items",
    "metrics": "KPIs, velocity, predictability, delivery metrics",
    "team": "Team or ART name",
    "date": "PI number or date"
}}

Return ONLY JSON. No markdown. No explanation.
"""
    try:
        response = model.generate_content(prompt)
        raw = response.text
        cleaned = clean_json_response(raw)
        return json.loads(cleaned)
    except Exception as e:
        raise Exception(f"LLM extraction failed: {e}")

def replace_text(shape, data):
    if not shape.has_text_frame:
        return
    for paragraph in shape.text_frame.paragraphs:
        full_text = ""
        for run in paragraph.runs:
            full_text += run.text
        updated_text = full_text
        for key, value in data.items():
            updated_text = updated_text.replace(key, str(value))
        if updated_text != full_text:
            paragraph.clear()
            paragraph.text = updated_text

def generate_ppt(template_path, data):
    if not os.path.exists(template_path):
        raise FileNotFoundError(f"Template not found: {template_path}")
    prs = Presentation(template_path)
    for slide in prs.slides:
        for shape in slide.shapes:
            replace_text(shape, data)
    os.makedirs("Output", exist_ok=True)
    output_path = "Output/generated_ppt.pptx"
    prs.save(output_path)
    return output_path